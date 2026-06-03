import os
import requests
"""import math"""
import chromadb
from pypdf import PdfReader
from openpyxl import load_workbook
from pptx import Presentation

#从文档文件夹中加载所有文本
def load_all_documents(folder_path):
    documents = []

    for filename in os.listdir(folder_path):
        path = os.path.join(folder_path, filename)

        if filename.startswith("."):
            continue

        if not filename.endswith((".pdf", ".xlsx", ".pptx", ".txt")):
            continue

        content = load_document(path)

        if filename.endswith(".pdf"):
            for page_data in content:
                documents.append(
                    {
                        "text": page_data["text"],
                        "source": filename,
                        "page": page_data["page"]
                    }
                )
        else:
            documents.append(
                {
                    "text": content,
                    "source": filename,
                    "page": None
                }
            )

    return documents

def load_pdf(path):
    reader = PdfReader(path)
    pages = []

    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            pages.append(
                {
                    "text": text,
                    "page": i + 1
                }
            )

    return pages

def load_excel(path):
    workbook = load_workbook(path, data_only=True)
    text = ""

    for sheet in workbook.worksheets:
        text += f"\nSheet: {sheet.title}\n"

        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell is not None])
            if row_text:
                text += row_text + "\n"

    return text

def load_ppt(path):
    presentation = Presentation(path)
    text = ""

    for i, slide in enumerate(presentation.slides):
        text += f"\nSlide {i + 1}\n"

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text

def load_txt(path):
    with open(path, "r", encoding="utf-8") as file:
        return file.read()

def load_document(path):
    if path.endswith(".pdf"):
        return load_pdf(path)

    if path.endswith(".xlsx"):
        return load_excel(path)

    if path.endswith(".pptx"):
        return load_ppt(path)

    if path.endswith(".txt"):
        return load_txt(path)

#将文本分割成较小的块，以便更好地处理和检索
def split_text(text):
    #简单地按句子分割文本
    sentences = text.replace("\n", " ").split(".")
    chunks = []

    #去除多余的空格并添加句号回到每个句子末尾
    for sentence in sentences:
        sentence = sentence.strip()
        if sentence:
            chunks.append(sentence + ".")

    return chunks

##获取文本的向量表示
def get_embedding(text):
    response = requests.post(
        "http://localhost:11434/api/embeddings",
        json={
            "model": "nomic-embed-text",
            "prompt": text
        }
    )

    data = response.json()
    return data["embedding"]

"""
======================

Legacy Retrieval
(Manual Similarity Search)

保留学习用途

======================

##计算两个向量之间的余弦相似度
def cosine_similarity(v1, v2):
    dot_product = sum(a * b for a, b in zip(v1, v2))
    magnitude1 = math.sqrt(sum(a * a for a in v1))
    magnitude2 = math.sqrt(sum(b * b for b in v2))

    return dot_product / (magnitude1 * magnitude2)

##查找与查询最相关的文本块
def retrieve_best_chunk(query, chunks):
    query_embedding = get_embedding(query)
    scores = []

    for chunk in chunks:
        chunk_embedding = get_embedding(chunk)
        similarity = cosine_similarity(query_embedding, chunk_embedding)
        scores.append((similarity, chunk))

    scores.sort(reverse=True)
    return scores[0][1]

"""

#构建向量数据库以便更高效地检索相关文本块
def build_vector_db(chunks):
    client = chromadb.Client()
    collection = client.get_or_create_collection(name="knowledgehub")

    for i, chunk_data in enumerate(chunks):
        chunk_text = chunk_data["text"]
        source_name = chunk_data["source"]
        page_number = chunk_data["page"]

        embedding = get_embedding(chunk_text)

        collection.add(
            ids=[str(i)],
            embeddings=[embedding],
            documents=[chunk_text],
            metadatas=[
                {
                    "source": source_name,
                    "page": page_number
                }
            ]
        )

    return collection

#从数据库中检索与查询最相关的文本块
def retrieve_from_db(query, collection):
    query_embedding = get_embedding(query)

    results = collection.query(
        query_embeddings=[query_embedding],
        n_results=2
    )

    return results["documents"][0], results["metadatas"][0]

def ask_llm(context, question):
    prompt = f"""
Answer the question based only on the context below.

Context:
{context}

Question:
{question}
"""

    response = requests.post(
        "http://localhost:11434/api/generate",
        json={
            "model": "qwen2.5:3b",
            "prompt": prompt,
            "stream": False
        }
    )

    data = response.json()
    return data["response"]

#主程序
documents = load_all_documents("documents")
all_chunks = []

for doc in documents:

    chunks = split_text(doc["text"])

    for chunk in chunks:

        all_chunks.append(
            {
                "text": chunk,
                "source": doc["source"],
                "page": doc["page"]
            }
        )

query = input("Ask a question: ")

collection = build_vector_db(all_chunks)

top_chunks, top_metadatas = retrieve_from_db(query, collection)

context = "\n".join(top_chunks)

answer = ask_llm(context, query)

print("\n===== Answer =====")
print(answer)

print("\n===== Sources =====")

for i, metadata in enumerate(top_metadatas):
    source = metadata["source"]
    page = metadata.get("page")

    if page:
        print(f"[{i+1}] {source} - Page {page}")
    else:
        print(f"[{i+1}] {source}")