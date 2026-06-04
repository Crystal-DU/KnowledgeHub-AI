import os
import requests
"""import math"""
import chromadb
from pypdf import PdfReader
from openpyxl import load_workbook
from pptx import Presentation

#从文档文件夹中加载所有文本
def load_all_documents(folder_path):
    documents = []

    for filename in os.listdir(folder_path):
        path = os.path.join(folder_path, filename)

        if filename.startswith("."):
            continue

        if not filename.endswith((".pdf", ".xlsx", ".pptx", ".txt")):
            continue

        content = load_document(path)

        if filename.endswith(".pdf"):
            for page_data in content:
                documents.append({
                    "text": page_data["text"],
                    "source": filename,
                    "page": page_data["page"],
                    "sheet": None,
                    "slide": None
                })

        elif filename.endswith(".xlsx"):
            for sheet_data in content:
                documents.append({
                    "text": sheet_data["text"],
                    "source": filename,
                    "page": None,
                    "sheet": sheet_data["sheet"],
                    "slide": None
                })

        elif filename.endswith(".pptx"):
            for slide_data in content:
                documents.append({
                    "text": slide_data["text"],
                    "source": filename,
                    "page": None,
                    "sheet": None,
                    "slide": slide_data["slide"]
                })

        elif filename.endswith(".txt"):
            documents.append({
                "text": content,
                "source": filename,
                "page": None,
                "sheet": None,
                "slide": None
            })

    return documents

def load_pdf(path):
    reader = PdfReader(path)
    pages = []

    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            pages.append(
                {
                    "text": text,
                    "page": i + 1
                }
            )

    return pages

def load_excel(path):
    workbook = load_workbook(path, data_only=True)
    sheets = []

    for sheet in workbook.worksheets:
        text = ""

        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell is not None])
            if row_text:
                text += row_text + "\n"

        if text:
            sheets.append(
                {
                    "text": text,
                    "sheet": sheet.title
                }
            )

    return sheets

def load_ppt(path):
    presentation = Presentation(path)
    slides = []

    for i, slide in enumerate(presentation.slides):
        text = ""

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

        if text:
            slides.append(
                {
                    "text": text,
                    "slide": i + 1
                }
            )

    return slides

def load_txt(path):
    with open(path, "r", encoding="utf-8") as file:
        return file.read()

def load_document(path):
    if path.endswith(".pdf"):
        return load_pdf(path)

    if path.endswith(".xlsx"):
        return load_excel(path)

    if path.endswith(".pptx"):
        return load_ppt(path)

    if path.endswith(".txt"):
        return load_txt(path)

#将文本分割成较小的块，以便更好地处理和检索
def split_text(text):
    #简单地按句子分割文本
    sentences = text.replace("\n", " ").split(".")
    chunks = []

    #去除多余的空格并添加句号回到每个句子末尾
    for sentence in sentences:
        sentence = sentence.strip()
        if sentence:
            chunks.append(sentence + ".")

    return chunks

##获取文本的向量表示
def get_embedding(text):
    response = requests.post(
        "http://localhost:11434/api/embeddings",
        json={
            "model": "nomic-embed-text",
            "prompt": text
        }
    )

    data = response.json()
    return data["embedding"]

"""
======================

Legacy Retrieval
(Manual Similarity Search)

保留学习用途

======================

##计算两个向量之间的余弦相似度
def cosine_similarity(v1, v2):
    dot_product = sum(a * b for a, b in zip(v1, v2))
    magnitude1 = math.sqrt(sum(a * a for a in v1))
    magnitude2 = math.sqrt(sum(b * b for b in v2))

    return dot_product / (magnitude1 * magnitude2)

##查找与查询最相关的文本块
def retrieve_best_chunk(query, chunks):
    query_embedding = get_embedding(query)
    scores = []

    for chunk in chunks:
        chunk_embedding = get_embedding(chunk)
        similarity = cosine_similarity(query_embedding, chunk_embedding)
        scores.append((similarity, chunk))

    scores.sort(reverse=True)
    return scores[0][1]

"""

#构建向量数据库以便更高效地检索相关文本块
def build_vector_db(chunks):
    client = chromadb.Client()
    collection = client.get_or_create_collection(name="knowledgehub")

    for i, chunk_data in enumerate(chunks):
        chunk_text = chunk_data["text"]
        source_name = chunk_data["source"]
        page_number = chunk_data["page"]
        sheet_name = chunk_data["sheet"]
        slide_number = chunk_data["slide"]

        embedding = get_embedding(chunk_text)
        metadata = {"source": source_name}
        if page_number is not None:
            metadata["page"] = page_number
        if sheet_name is not None:
            metadata["sheet"] = sheet_name
        if slide_number is not None:
            metadata["slide"] = slide_number

        collection.add(
            ids=[str(i)],
            embeddings=[embedding],
            documents=[chunk_text],
            metadatas=[metadata]
        )
    return collection

#从数据库中检索与查询最相关的文本块
def retrieve_from_db(query, collection):
    query_embedding = get_embedding(query)

    results = collection.query(
        query_embeddings=[query_embedding],
        n_results=2
    )

    return results["documents"][0], results["metadatas"][0]

def ask_llm(context, question):
    prompt = f"""
Answer the question based only on the context below.

Context:
{context}

Question:
{question}
"""

    response = requests.post(
        "http://localhost:11434/api/generate",
        json={
            "model": "qwen2.5:3b",
            "prompt": prompt,
            "stream": False
        }
    )

    data = response.json()
    return data["response"]

def ask_question(query):
    # 1. 从 documents 文件夹读取所有文档
    documents = load_all_documents("documents")

    # 2. 准备一个列表，用来保存所有 chunk
    all_chunks = []

    # 3. 遍历每个文档
    for doc in documents:

        # 4. 把这个文档的 text 切成多个 chunk
        split_chunks = split_text(doc["text"])

        # 5. 给每个 chunk 加上来源信息
        for chunk in split_chunks:
            all_chunks.append(
                {
                    "text": chunk,
                    "source": doc["source"],
                    "page": doc.get("page"),
                    "sheet": doc.get("sheet"),
                    "slide": doc.get("slide")
                }
            )

    # 6. 把所有 chunk 存进 ChromaDB
    collection = build_vector_db(all_chunks)

    # 7. 根据用户问题检索最相关的 chunk
    top_chunks, top_metadatas = retrieve_from_db(
        query,
        collection
    )

    # 8. 把检索到的 chunk 合并成 context
    context = "\n".join(top_chunks)

    # 9. 让 Qwen 根据 context 回答
    answer = ask_llm(context, query)

    # 10. 返回答案和来源
    return answer, top_metadatas